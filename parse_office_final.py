#!/usr/bin/env python3
"""Final parser for office documents with proper multi-paragraph handling"""
import json
import re
from pathlib import Path
from docx import Document
from pptx import Presentation
import hashlib

def normalize_for_hash(text):
    text = text.lower().strip()
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^\w\s]', '', text)
    return text

def extract_from_docx(filepath):
    """Extract questions from docx where Q and options span multiple paragraphs"""
    questions = []
    
    try:
        doc = Document(filepath)
        paragraphs = [p.text.strip() for p in doc.paragraphs]
        
        i = 0
        while i < len(paragraphs):
            # Look for "Question N" marker
            if re.match(r'^Question\s+\d+', paragraphs[i], re.IGNORECASE):
                question_parts = []
                options = {}
                
                # Move past the "Question N" line
                i += 1
                
                # Collect question text (until we hit an option or another Question)
                while i < len(paragraphs):
                    line = paragraphs[i]
                    
                    if not line:  # Empty line
                        i += 1
                        continue
                    
                    # Check if this is an option
                    if re.match(r'^([A-D])\s+', line):
                        match = re.match(r'^([A-D])\s+(.*)', line)
                        if match:
                            options[match.group(1)] = match.group(2).strip()
                        i += 1
                    # Check if it's the next question
                    elif re.match(r'^Question\s+\d+', line, re.IGNORECASE):
                        break
                    # Otherwise it's part of the question or continuation of option
                    else:
                        # If we have options already, this might be continuation of last option
                        if options and not '?' in line:
                            last_key = sorted(options.keys())[-1]
                            options[last_key] += ' ' + line
                        else:
                            # Part of question
                            question_parts.append(line)
                        i += 1
                
                # Combine question parts
                question_text = ' '.join(question_parts).strip()
                
                # Validate and add
                if question_text and len(options) >= 2 and '?' in question_text:
                    questions.append({
                        'question': question_text,
                        'options': options,
                        'source': filepath.name
                    })
            else:
                i += 1
                
    except Exception as e:
        print(f"Error in {filepath.name}: {e}")
    
    return questions

def extract_from_pptx(filepath):
    """Extract questions from PPTX slides"""
    questions = []
    
    try:
        prs = Presentation(filepath)
        
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            
            if not texts:
                continue
            
            # Try to identify question and options
            question = None
            options = {}
            
            for text in texts:
                text = text.strip()
                
                # Check if it's a question
                if '?' in text and not question:
                    # Remove "Question N" prefix if present
                    text = re.sub(r'^Question\s+\d+[\s\n]+', '', text, flags=re.IGNORECASE|re.MULTILINE)
                    question = text.strip()
                
                # Check if it's an option
                if re.match(r'^([A-D])[\.\):\s]+', text):
                    match = re.match(r'^([A-D])[\.\):\s]+(.+)', text, re.DOTALL)
                    if match:
                        key = match.group(1)
                        value = match.group(2).strip()
                        options[key] = value
            
            # Validate
            if question and len(options) >= 2 and len(question) >= 20:
                questions.append({
                    'question': question,
                    'options': options,
                    'source': filepath.name
                })
                
    except Exception as e:
        print(f"Error in {filepath.name}: {e}")
    
    return questions

def detect_language(text):
    """Simple language detection"""
    french_words = ['le', 'la', 'les', 'un', 'une', 'des', 'est', 'sont', 'quel', 
                    'quelle', 'comment', 'dans', 'qui', 'avec', 'pour', 'par']
    english_words = ['the', 'is', 'are', 'what', 'which', 'how', 'that', 'this', 
                     'shall', 'will', 'would', 'should', 'under', 'with', 'from']
    
    text_lower = text.lower()
    fr_count = sum(1 for w in french_words if f' {w} ' in f' {text_lower} ')
    en_count = sum(1 for w in english_words if f' {w} ' in f' {text_lower} ')
    
    return 'fr' if fr_count > en_count else 'en'

# Main processing
print("="*70)
print("PARSING OFFICE FILES")
print("="*70)

# Load existing dataset
existing = json.load(open('quiz-data-translated.json', encoding='utf-8'))
existing_hashes = set()

for q in existing:
    for lang in ['en', 'fr', 'ru']:
        qtxt = q.get('translations', {}).get(lang, {}).get('question', '')
        if qtxt:
            existing_hashes.add(hashlib.md5(normalize_for_hash(qtxt).encode()).hexdigest())

print(f"Loaded {len(existing)} existing questions")
print()

# Extract from all files
src_dir = Path('source/SRC')
all_questions = []

print("Extracting from DOCX files:")
for docx_file in sorted(src_dir.glob('*.docx')):
    if not docx_file.name.startswith('.~lock'):
        questions = extract_from_docx(docx_file)
        all_questions.extend(questions)
        print(f"  {docx_file.name:50} {len(questions):3} questions")

print("\nExtracting from PPTX files:")
for pptx_file in sorted(src_dir.glob('*.pptx')):
    if not pptx_file.name.startswith('.~lock'):
        questions = extract_from_pptx(pptx_file)
        all_questions.extend(questions)
        print(f"  {pptx_file.name:50} {len(questions):3} questions")

print(f"\nTotal questions extracted: {len(all_questions)}")

# Deduplicate
seen_hashes = set()
new_questions = []

for item in all_questions:
    qhash = hashlib.md5(normalize_for_hash(item['question']).encode()).hexdigest()
    
    if qhash not in existing_hashes and qhash not in seen_hashes:
        lang = detect_language(item['question'])
        
        structured = {
            'id': None,
            'correctOptionKey': 'B',  # Default - needs manual review
            'translations': {
                lang: {
                    'question': item['question'],
                    'options': item['options']
                }
            },
            'images': [],
            'source': item['source']
        }
        
        new_questions.append(structured)
        seen_hashes.add(qhash)

# Assign IDs
next_id = max(q['id'] for q in existing) + 1
for i, q in enumerate(new_questions):
    q['id'] = next_id + i

print(f"\n{'='*70}")
print(f"RESULTS")
print(f"{'='*70}")
print(f"New unique questions: {len(new_questions)}")

# Language stats
lang_stats = {'en': 0, 'fr': 0}
for q in new_questions:
    for lang in q['translations']:
        lang_stats[lang] += 1

print(f"  English: {lang_stats['en']}")
print(f"  French: {lang_stats['fr']}")

# Save
output_file = 'quiz-data-office-new.json'
with open(output_file, 'w', encoding='utf-8') as f:
    json.dump(new_questions, f, ensure_ascii=False, indent=2)

print(f"\nSaved to {output_file}")

# Show samples
if lang_stats['en'] > 0:
    print(f"\n{'='*70}")
    print("ENGLISH SAMPLES (first 3):")
    print(f"{'='*70}")
    en_questions = [q for q in new_questions if 'en' in q['translations']]
    for q in en_questions[:3]:
        trans = q['translations']['en']
        print(f"\nQ{q['id']} [{q['source']}]:")
        print(f"  {trans['question'][:100]}...")
        for opt_key in sorted(trans['options'].keys()):
            print(f"    {opt_key}. {trans['options'][opt_key][:70]}...")

if lang_stats['fr'] > 0:
    print(f"\n{'='*70}")
    print("FRENCH SAMPLES (first 3):")
    print(f"{'='*70}")
    fr_questions = [q for q in new_questions if 'fr' in q['translations']]
    for q in fr_questions[:3]:
        trans = q['translations']['fr']
        print(f"\nQ{q['id']} [{q['source']}]:")
        print(f"  {trans['question'][:100]}...")
        for opt_key in sorted(trans['options'].keys()):
            print(f"    {opt_key}. {trans['options'][opt_key][:70]}...")

print(f"\n{'='*70}")
print("NEXT STEPS:")
print(f"{'='*70}")
print("1. Review quiz-data-office-new.json for quality")
print("2. Translate missing language versions (EN→FR or FR→EN)")
print("3. Merge with quiz-data-translated.json")
print("4. Review and fix correctOptionKey values")
print("5. Generate final quiz-data.js for React app")
