#!/usr/bin/env python3
"""
Parse quiz questions from docx/pptx files
"""
import json
import re
from pathlib import Path
from docx import Document
from pptx import Presentation
import hashlib

def normalize_for_hash(text):
    """Normalize text for deduplication"""
    text = text.lower().strip()
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^\w\s]', '', text)
    return text

def extract_questions_from_docx(filepath):
    """Extract structured questions from docx files"""
    raw_questions = []
    
    try:
        doc = Document(filepath)
        current_q = {
            'question': '',
            'options': {},
            'lines': []
        }
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                # Empty line - might be end of question
                if current_q['question'] and current_q['options']:
                    raw_questions.append(current_q)
                    current_q = {'question': '', 'options': {}, 'lines': []}
                continue
            
            # Check if this is a question number or new question
            if re.match(r'^(Question\s+\d+|Q\d+|\d+[\.\)])\s+', text, re.IGNORECASE):
                if current_q['question'] and current_q['options']:
                    raw_questions.append(current_q)
                # Extract question without the number
                question_text = re.sub(r'^(Question\s+\d+|Q\d+|\d+[\.\)])\s+', '', text, flags=re.IGNORECASE)
                current_q = {
                    'question': question_text,
                    'options': {},
                    'lines': [question_text]
                }
            # Check if this is an option (A., B., C., D.)
            elif re.match(r'^([A-D])[\.:\)]\s+.+', text):
                match = re.match(r'^([A-D])[\.:\)]\s+(.+)', text)
                if match:
                    key = match.group(1)
                    value = match.group(2).strip()
                    if value:  # Only add non-empty options
                        current_q['options'][key] = value
                        current_q['lines'].append(text)
            # Check if it's a question line (has ?)
            elif '?' in text:
                if not current_q['question']:
                    current_q['question'] = text
                    current_q['lines'] = [text]
                else:
                    # Append to existing question
                    current_q['question'] += ' ' + text
                    current_q['lines'].append(text)
            # Otherwise might be continuation of question
            elif current_q['question'] and not current_q['options']:
                current_q['question'] += ' ' + text
                current_q['lines'].append(text)
        
        # Don't forget the last question
        if current_q['question'] and current_q['options']:
            raw_questions.append(current_q)
            
    except Exception as err:
        print(f"Error parsing {filepath.name}: {err}")
    
    # Filter out invalid questions
    valid_questions = []
    for q in raw_questions:
        # Must have a question with at least 20 chars, must have ?, and at least 2 options
        if (len(q['question']) >= 20 and 
            '?' in q['question'] and 
            len(q['options']) >= 2 and
            not q['question'].startswith('SAMPLE QUESTIONS')):
            valid_questions.append(q)
    
    return valid_questions

def extract_questions_from_pptx(filepath):
    """Extract structured questions from pptx files"""
    questions = []
    
    try:
        prs = Presentation(filepath)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if not slide_text:
                continue
            
            # Try to identify question structure
            question_text = ""
            options = {}
            
            for text in slide_text:
                # Check if this looks like a question
                if '?' in text and not question_text:
                    question_text = text
                # Check for options
                elif re.match(r'^([A-D])[\.:\)]', text):
                    match = re.match(r'^([A-D])[\.:\)]\s*(.*)', text)
                    if match:
                        options[match.group(1)] = match.group(2)
            
            # If we found a question, add it
            if question_text:
                questions.append({
                    'question': question_text,
                    'options': options,
                    'lines': slide_text
                })
            # Otherwise, if there's content that looks like it might be a question
            elif slide_text and any(word in ' '.join(slide_text).lower() for word in ['what', 'which', 'how', 'when', 'where', 'quel', 'quelle', 'comment', 'quand']):
                combined = '\n'.join(slide_text)
                questions.append({
                    'question': combined.split('\n')[0],
                    'options': {},
                    'lines': slide_text
                })
    
    except Exception as e:
        print(f"Error parsing {filepath.name}: {e}")
    
    return questions

def detect_language(text):
    """Simple language detection based on common words"""
    french_indicators = ['le', 'la', 'les', 'un', 'une', 'des', 'est', 'sont', 'quel', 'quelle', 'quels']
    english_indicators = ['the', 'is', 'are', 'what', 'which', 'how', 'that', 'this']
    
    text_lower = text.lower()
    french_count = sum(1 for word in french_indicators if f' {word} ' in f' {text_lower} ')
    english_count = sum(1 for word in english_indicators if f' {word} ' in f' {text_lower} ')
    
    if french_count > english_count:
        return 'fr'
    return 'en'

# Load existing questions to check for duplicates
existing_data = json.load(open('quiz-data-translated.json', encoding='utf-8'))
existing_hashes = set()
for q in existing_data:
    for lang in ['en', 'fr', 'ru']:
        qtxt = q.get('translations', {}).get(lang, {}).get('question', '')
        if qtxt:
            existing_hashes.add(hashlib.md5(normalize_for_hash(qtxt).encode()).hexdigest())

# Parse all docx and pptx files
src_dir = Path('source/SRC')
all_questions = []

print("Parsing DOCX files...")
for docx_file in sorted(src_dir.glob('*.docx')):
    if not docx_file.name.startswith('.~lock'):
        print(f"  - {docx_file.name}")
        questions = extract_questions_from_docx(docx_file)
        for q in questions:
            q['source'] = docx_file.name
        all_questions.extend(questions)

print("\nParsing PPTX files...")
for pptx_file in sorted(src_dir.glob('*.pptx')):
    if not pptx_file.name.startswith('.~lock'):
        print(f"  - {pptx_file.name}")
        questions = extract_questions_from_pptx(pptx_file)
        for q in questions:
            q['source'] = pptx_file.name
        all_questions.extend(questions)

print(f"\nTotal questions extracted: {len(all_questions)}")

# Filter out duplicates against existing dataset
new_questions = []
for q in all_questions:
    qhash = hashlib.md5(normalize_for_hash(q['question']).encode()).hexdigest()
    if qhash not in existing_hashes:
        # Structure the question data
        lang = detect_language(q['question'])
        
        structured = {
            'id': None,  # Will assign later
            'correctOptionKey': 'B',  # Default - needs manual review
            'translations': {},
            'images': [],
            'source': q['source']
        }
        
        structured['translations'][lang] = {
            'question': q['question'],
            'options': q['options'] if q['options'] else {
                'A': '(option A - needs extraction)',
                'B': '(option B - needs extraction)',
                'C': '(option C - needs extraction)',
                'D': '(option D - needs extraction)'
            }
        }
        
        new_questions.append(structured)
        existing_hashes.add(qhash)  # Prevent duplicates within this batch

print(f"New unique questions (not in existing dataset): {len(new_questions)}")

# Assign IDs starting after the existing ones
next_id = max(q['id'] for q in existing_data) + 1
for i, q in enumerate(new_questions):
    q['id'] = next_id + i

# Save the new questions
output_file = 'quiz-data-office-new.json'
with open(output_file, 'w', encoding='utf-8') as f:
    json.dump(new_questions, f, ensure_ascii=False, indent=2)

print(f"\nSaved {len(new_questions)} new questions to {output_file}")

# Show statistics
lang_stats = {'en': 0, 'fr': 0}
with_options = 0
for q in new_questions:
    for lang in q['translations']:
        lang_stats[lang] = lang_stats.get(lang, 0) + 1
        if len(q['translations'][lang].get('options', {})) >= 4:
            with_options += 1

print(f"\nStatistics:")
print(f"  English questions: {lang_stats.get('en', 0)}")
print(f"  French questions: {lang_stats.get('fr', 0)}")
print(f"  Questions with 4+ options: {with_options}")
print(f"  Questions needing option extraction: {len(new_questions) - with_options}")

# Show samples
print(f"\nSample new questions (first 5):")
for i, q in enumerate(new_questions[:5], 1):
    lang = list(q['translations'].keys())[0]
    qtext = q['translations'][lang]['question']
    opts = q['translations'][lang]['options']
    print(f"\n{i}. [{lang.upper()}] {qtext[:80]}...")
    print(f"   Source: {q['source']}")
    print(f"   Options: {len(opts)} defined")
